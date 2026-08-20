import os
from typing import Optional
import docx
from pptx import Presentation
from app.services.ollama_translator import call_ollama_generate, OLLAMA_DEFAULT_MODEL
from app.utils.logger import logger

def translate_single_text(
    text: str,
    source_lang: str = "en",
    target_lang: str = "vi",
    model: str = OLLAMA_DEFAULT_MODEL
) -> str:
    """Helper dịch một chuỗi văn bản ngắn/đoạn qua Ollama."""
    clean_text = text.strip()
    if not clean_text or len(clean_text) <= 3:
        return text

    prompt = (
        f"Dịch ngắn gọn đoạn văn bản sau từ tiếng {source_lang.upper()} sang tiếng {target_lang.upper()}.\n"
        f"Chỉ trả về bản dịch, không giải thích thừa.\n\n"
        f"Text:\n{clean_text}"
    )
    try:
        translated = call_ollama_generate(prompt=prompt, model=model)
        return translated if translated else text
    except Exception as e:
        logger.warning(f"Lỗi dịch đoạn ngắn bằng Ollama ({e}), thử dùng Gemini API fallback...")
        try:
            from app.services.translator import translate_chunk_with_gemini
            return translate_chunk_with_gemini(clean_text, source_lang=source_lang, target_lang=target_lang)
        except Exception as gemini_err:
            logger.error(f"Lỗi cả Gemini API fallback: {gemini_err}")
            return text

from typing import Optional, Tuple

def translate_docx_document(
    input_path: str,
    output_path: str,
    source_lang: str = "en",
    target_lang: str = "vi",
    model: str = OLLAMA_DEFAULT_MODEL
) -> Tuple[str, str, str]:
    """
    Dịch file Microsoft Word (.docx) In-place:
    - Duyệt qua doc.paragraphs, dịch các đoạn > 5 ký tự và gán lại para.text.
    - Duyệt qua doc.tables, từng row, cell. Dịch cell.text và gán lại.
    - Trả về: (output_path, translated_markdown_text, raw_source_text)
    """
    logger.info(f"Bắt đầu dịch file Word .docx: {input_path} -> {output_path}")
    doc = docx.Document(input_path)
    
    translated_lines = []
    source_lines = []

    # 1. Dịch các đoạn văn (paragraphs)
    for para in doc.paragraphs:
        raw = para.text.strip()
        if raw and len(raw) > 3:
            source_lines.append(raw)
            translated_text = translate_single_text(
                raw, source_lang=source_lang, target_lang=target_lang, model=model
            )
            para.text = translated_text
            translated_lines.append(translated_text)

    # 2. Dịch các bảng biểu (tables) giữ nguyên cấu trúc khung bảng (borders & layout)
    for table in doc.tables:
        for row in table.rows:
            row_translated = []
            for cell in row.cells:
                for para in cell.paragraphs:
                    raw = para.text.strip()
                    if raw and len(raw) > 1:
                        source_lines.append(raw)
                        translated_para = translate_single_text(
                            raw, source_lang=source_lang, target_lang=target_lang, model=model
                        )
                        para.text = translated_para
                        row_translated.append(translated_para)
            if row_translated:
                translated_lines.append(" | ".join(row_translated))

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    doc.save(output_path)
    logger.info(f"Đã hoàn thành dịch file Word .docx và lưu tại {output_path}")
    
    translated_md = "\n\n".join(translated_lines)
    source_md = "\n\n".join(source_lines)
    return output_path, translated_md, source_md

def translate_pptx_document(
    input_path: str,
    output_path: str,
    source_lang: str = "en",
    target_lang: str = "vi",
    model: str = OLLAMA_DEFAULT_MODEL
) -> Tuple[str, str, str]:
    """
    Dịch file Microsoft PowerPoint (.pptx) In-place:
    - Duyệt qua các slide -> shapes -> text_frame.
    - Dịch paragraph.text và gán lại.
    - Bật thuộc tính AutoFit (word_wrap = True) để chữ tiếng Việt dài không tràn khỏi box.
    - Trả về: (output_path, translated_markdown_text, raw_source_text)
    """
    logger.info(f"Bắt đầu dịch file PowerPoint .pptx: {input_path} -> {output_path}")
    prs = Presentation(input_path)

    translated_slides = []
    source_lines = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                text_frame.word_wrap = True # AutoFit / word wrap
                
                for para in text_frame.paragraphs:
                    raw = para.text.strip()
                    if raw and len(raw) > 2:
                        source_lines.append(raw)
                        translated_text = translate_single_text(
                            raw, source_lang=source_lang, target_lang=target_lang, model=model
                        )
                        para.text = translated_text
                        slide_texts.append(f"- {translated_text.strip()}")
        
        if slide_texts:
            translated_slides.append(f"### 🖥️ Slide {slide_idx}\n" + "\n".join(slide_texts))

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    logger.info(f"Đã hoàn thành dịch file PowerPoint .pptx và lưu tại {output_path}")
    
    translated_md = "\n\n".join(translated_slides)
    source_md = "\n\n".join(source_lines)
    return output_path, translated_md, source_md
